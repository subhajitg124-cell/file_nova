import os
import shutil
import subprocess
from pathlib import Path
from typing import List, Dict, Any, Callable, Optional
from openpyxl import load_workbook
from docx import Document
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from app.processors.base import BaseProcessor
from app.config import settings

class DocumentProcessor(BaseProcessor):
    def validate_options(self, options: Dict[str, Any]) -> None:
        operation = options.get("operation")
        if not operation:
            raise ValueError("Missing 'operation' in options.")
        if operation not in ["docx_to_pdf", "pptx_to_pdf", "xlsx_to_csv", "docx_cleanup", "md_to_html", "docx_merge"]:
            raise ValueError(f"Unsupported Document operation: {operation}")

    def process(
        self, 
        file_paths: List[str], 
        options: Dict[str, Any], 
        progress_callback: Optional[Callable[[float], None]] = None
    ) -> str:
        self.validate_options(options)
        operation = options["operation"]
        
        if not file_paths:
            raise ValueError("No files provided for document processing.")
            
        input_path = file_paths[0]
        if not os.path.exists(input_path):
            raise FileNotFoundError(f"Input document not found: {input_path}")
            
        if progress_callback:
            progress_callback(10.0)
            
        output_path = None
        
        if operation == "docx_to_pdf":
            output_path = self._convert_to_pdf_libreoffice_or_fallback(input_path, progress_callback)
        elif operation == "pptx_to_pdf":
            output_path = self._convert_to_pdf_libreoffice_or_fallback(input_path, progress_callback)
        elif operation == "xlsx_to_csv":
            output_path = self._xlsx_to_csv(input_path, progress_callback)
        elif operation == "docx_cleanup":
            output_path = self._docx_cleanup(input_path, progress_callback)
        elif operation == "md_to_html":
            output_path = self._md_to_html(input_path, progress_callback)
        elif operation == "docx_merge":
            if not file_paths or len(file_paths) < 2:
                raise ValueError("Merging requires at least 2 files.")
            output_path = self._merge_docx(file_paths, progress_callback)
            
        if progress_callback:
            progress_callback(100.0)
            
        return str(output_path)

    def _get_libreoffice_executable(self) -> Optional[str]:
        """Checks if LibreOffice is available in the system PATH."""
        for cmd in ["soffice", "libreoffice", "soffice.exe"]:
            if shutil.which(cmd):
                return cmd
        # Common installation paths on Windows
        win_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"
        ]
        for p in win_paths:
            if os.path.exists(p):
                return p
        return None

    def _convert_to_pdf_libreoffice_or_fallback(
        self, 
        input_path: str, 
        progress_callback: Optional[Callable[[float], None]]
    ) -> Path:
        output_filename = f"{Path(input_path).stem}.pdf"
        output_file_path = settings.output_dir / output_filename
        
        libreoffice_bin = self._get_libreoffice_executable()
        
        if libreoffice_bin:
            if progress_callback:
                progress_callback(30.0)
            
            # Execute LibreOffice Headless PDF conversion
            try:
                cmd = [
                    libreoffice_bin,
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(settings.output_dir),
                    input_path
                ]
                subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                
                if progress_callback:
                    progress_callback(80.0)
                    
                # The output file is named [stem].pdf in the outdir
                created_pdf = settings.output_dir / f"{Path(input_path).stem}.pdf"
                if created_pdf.exists():
                    return created_pdf
            except Exception as e:
                # If libreoffice fails, log it and run fallback
                pass
                
        # FALLBACK: Create a PDF using ReportLab with the text content of the DOCX or PPTX
        if progress_callback:
            progress_callback(40.0)
            
        doc = SimpleDocTemplate(str(output_file_path), pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        # Add Header
        story.append(Paragraph(f"<b>Document Preview Fallback:</b> {Path(input_path).name}", styles['Title']))
        story.append(Spacer(1, 15))
        story.append(Paragraph("<i>Note: This document was rendered via ReportLab fallback because headless LibreOffice was not available on the server.</i>", styles['Italic']))
        story.append(Spacer(1, 20))
        
        # Read text from file
        content_extracted = ""
        if input_path.endswith(".docx"):
            try:
                docx_doc = Document(input_path)
                content_extracted = "\n".join([p.text for p in docx_doc.paragraphs if p.text.strip()])
            except Exception:
                content_extracted = "Failed to extract text from DOCX document."
        elif input_path.endswith(".pptx"):
            try:
                from pptx import Presentation
                prs = Presentation(input_path)
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_runs.append(shape.text)
                content_extracted = "\n\n".join(text_runs)
            except Exception:
                content_extracted = "Failed to extract text from PPTX presentation."
        else:
            content_extracted = "Unsupported document extension for text fallback."

        # Add text content to pdf
        for paragraph_text in content_extracted.split("\n"):
            if paragraph_text.strip():
                story.append(Paragraph(paragraph_text, styles['Normal']))
                story.append(Spacer(1, 8))
                
        doc.build(story)
        
        if progress_callback:
            progress_callback(90.0)
            
        return output_file_path

    def _xlsx_to_csv(self, input_path: str, progress_callback: Optional[Callable[[float], None]]) -> Path:
        output_filename = f"{Path(input_path).stem}.csv"
        output_file_path = settings.output_dir / output_filename
        
        wb = load_workbook(input_path, read_only=True, data_only=True)
        sheet = wb.active
        
        import csv
        with open(output_file_path, "w", newline="", encoding="utf-8") as f_out:
            writer = csv.writer(f_out)
            for row in sheet.iter_rows(values_only=True):
                # Write rows to CSV
                writer.writerow(list(row))
                
        if progress_callback:
            progress_callback(80.0)
            
        return output_file_path

    def _docx_cleanup(self, input_path: str, progress_callback: Optional[Callable[[float], None]]) -> Path:
        output_filename = f"cleaned_{Path(input_path).name}"
        output_file_path = settings.output_dir / output_filename
        
        doc = Document(input_path)
        
        # Enforce standard formatting
        # 1. Remove empty paragraphs
        paragraphs_to_remove = []
        for p in doc.paragraphs:
            if not p.text.strip():
                paragraphs_to_remove.append(p)
        
        # Note: XML manipulation is needed to safely remove a paragraph in python-docx
        for p in paragraphs_to_remove:
            p_element = p._element
            p_element.getparent().remove(p_element)
            
        if progress_callback:
            progress_callback(50.0)
            
        # 2. Normalize margins (e.g. set standard 1 inch margins)
        for section in doc.sections:
            section.top_margin = Inches = 72 * 20 * 20  # Equivalent to docx.shared.Inches(1) in XML EMU
            section.bottom_margin = Inches
            section.left_margin = Inches
            section.right_margin = Inches
            
        doc.save(str(output_file_path))
        
        if progress_callback:
            progress_callback(85.0)
            
        return output_file_path

    def _md_to_html(self, input_path: str, progress_callback: Optional[Callable[[float], None]]) -> Path:
        output_filename = f"{Path(input_path).stem}.html"
        output_file_path = settings.output_dir / output_filename
        
        with open(input_path, "r", encoding="utf-8") as f_in:
            md_text = f_in.read()
            
        # Basic markdown to HTML renderer (supporting headers, lists, links, bold, italics)
        html_lines = []
        html_lines.append("<!DOCTYPE html>")
        html_lines.append("<html>")
        html_lines.append("<head><style>body { font-family: sans-serif; line-height: 1.6; padding: 20px; max-width: 800px; margin: 0 auto; }</style></head>")
        html_lines.append("<body>")
        
        in_list = False
        for line in md_text.split("\n"):
            line = line.strip()
            if not line:
                if in_list:
                    html_lines.append("</ul>")
                    in_list = False
                continue
                
            # Headers
            if line.startswith("# "):
                html_lines.append(f"<h1>{line[2:]}</h1>")
            elif line.startswith("## "):
                html_lines.append(f"<h2>{line[3:]}</h2>")
            elif line.startswith("### "):
                html_lines.append(f"<h3>{line[4:]}</h3>")
            # Lists
            elif line.startswith("- ") or line.startswith("* "):
                if not in_list:
                    html_lines.append("<ul>")
                    in_list = True
                html_lines.append(f"<li>{line[2:]}</li>")
            else:
                if in_list:
                    html_lines.append("</ul>")
                    in_list = False
                # Paragraph
                # basic bold / italic parse
                parsed_line = line.replace("**", "<b>", 1).replace("**", "</b>", 1)
                parsed_line = parsed_line.replace("*", "<i>", 1).replace("*", "</i>", 1)
                html_lines.append(f"<p>{parsed_line}</p>")
                
        if in_list:
            html_lines.append("</ul>")
            
        html_lines.append("</body>")
        html_lines.append("</html>")
        
        with open(output_file_path, "w", encoding="utf-8") as f_out:
            f_out.write("\n".join(html_lines))
            
        if progress_callback:
            progress_callback(80.0)
            
        return output_file_path

    def _merge_docx(
        self, 
        file_paths: List[str], 
        progress_callback: Optional[Callable[[float], None]]
    ) -> Path:
        output_filename = f"merged_{Path(file_paths[0]).name}"
        output_file_path = settings.output_dir / output_filename
        
        base_doc = Document(file_paths[0])
        total_files = len(file_paths)
        
        for idx, file_path in enumerate(file_paths[1:]):
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"Input DOCX file not found: {file_path}")
                
            # Add page break before appending content from next document
            base_doc.add_page_break()
            
            doc_to_append = Document(file_path)
            body = base_doc.element.body
            sectPr = body.sectPr
            
            if sectPr is not None:
                sectPr_index = body.index(sectPr)
            else:
                sectPr_index = len(body)
                
            for element in doc_to_append.element.body:
                if element.tag.endswith('sectPr'):
                    continue
                if sectPr is not None:
                    body.insert(sectPr_index, element)
                    sectPr_index += 1
                else:
                    body.append(element)
                    
            if progress_callback:
                # scale progress from 10% to 90%
                current_prog = 10.0 + ((idx + 1) / (total_files - 1)) * 80.0
                progress_callback(min(current_prog, 90.0))
                
        base_doc.save(str(output_file_path))
        
        if progress_callback:
            progress_callback(95.0)
            
        return output_file_path
